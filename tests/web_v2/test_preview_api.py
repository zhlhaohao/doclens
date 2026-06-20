"""GET /api/preview 测试。"""
import asyncio

import pytest
from httpx import ASGITransport, AsyncClient

from cortex.web_v2 import deps
from cortex.web_v2.app import create_app


@pytest.fixture
def reset_deps():
    """每个测试前后重置 deps 单例，避免 search_path 跨测试污染。"""
    deps.reset_singletons()
    yield
    deps.reset_singletons()


@pytest.mark.asyncio
async def test_preview_returns_file_content(temp_workdir, env_cortex_config, reset_deps):
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "doc1.md"})
    assert res.status_code == 200
    body = res.json()
    assert "Hello world" in body["content"]
    assert body["path"] == "doc1.md"


@pytest.mark.asyncio
async def test_preview_404_for_missing_file(temp_workdir, env_cortex_config, reset_deps):
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "nonexistent.md"})
    assert res.status_code == 404
    assert res.json()["code"] == "FILE_NOT_FOUND"


# ---------------------------------------------------------------------------
# Binary preview (Task 4) 集成测试
# ---------------------------------------------------------------------------


def _init_and_reindex():
    """在非事件循环线程中初始化 IndexManager 并建索引。

    TreeSearch.index 是同步阻塞调用，不能直接在 pytest-asyncio 事件循环里运行，
    否则会卡住（SQLite lock / thread affinity）。调用方需用 asyncio.to_thread 包装。
    """
    idx = deps.get_index_manager()
    idx.reindex(force=True)
    return idx


@pytest.mark.asyncio
async def test_preview_csv_returns_markdown(temp_workdir, env_cortex_config, reset_deps):
    """csv 预览应从 DB 合成 markdown，渲染为 md table。"""
    # 先触发索引，让 data.csv 进入 DB
    await asyncio.to_thread(_init_and_reindex)

    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "data.csv"})
    assert res.status_code == 200
    body = res.json()
    assert body["language"] == "markdown"
    assert "# data" in body["content"]
    # 应渲染为 md table
    assert "| name | age |" in body["content"]
    assert "| Alice | 30 |" in body["content"]


@pytest.mark.asyncio
async def test_preview_pptx_returns_markdown_from_db(
    temp_workdir, env_cortex_config, reset_deps
):
    """pptx 是二进制（ZIP），预览应从 DB 合成 markdown，而非 utf-8 直读（会乱码）。

    Regression: 之前 .pptx 没在 BINARY_PREVIEW_EXTS 里，导致走 utf-8 分支，
    把 ZIP 字节流当文本解码 → 乱码。修复后应通过 markitdown 入库的结构合成 md。
    """
    # 用 python-pptx 造一个最小 pptx（markitdown 依赖 python-pptx，可用）
    from pptx import Presentation  # type: ignore

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX Title"
    slide.placeholders[1].text = "Slide body content"

    pptx_path = temp_workdir / "sample.pptx"
    prs.save(str(pptx_path))
    assert pptx_path.exists()

    # 索引（让 markitdown parser 把 pptx 写入 DB）
    await asyncio.to_thread(_init_and_reindex)

    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "sample.pptx"})

    assert res.status_code == 200
    body = res.json()
    # 关键：language 必须是 markdown（走 DB 合成），不能是 text（utf-8 直读）
    assert body["language"] == "markdown", (
        f"pptx 应走二进制合成路径，实际 language={body['language']!r}, "
        f"content 前 80 字节: {body['content'][:80]!r}"
    )
    # markitdown 应提取出标题
    assert "Hello PPTX Title" in body["content"]
    # 不能出现 ZIP magic（PK\x03\x04）—— 那是 utf-8 直读乱码的标志
    assert "PK\x03\x04" not in body["content"]


@pytest.mark.asyncio
async def test_preview_unindexed_pdf_returns_404(temp_workdir, env_cortex_config, reset_deps):
    """未索引的 .pdf 文件应返回 404 NOT_INDEXED（不读磁盘 utf-8）。"""
    # 故意不索引；temp_workdir 里也没有这个文件
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "missing.pdf"})
    assert res.status_code == 404
    body = res.json()
    assert body["code"] == "NOT_INDEXED"


@pytest.mark.asyncio
async def test_preview_txt_still_uses_text_route(temp_workdir, env_cortex_config, reset_deps):
    """回归：.txt 仍走 utf-8 纯文本路径，language='text'。"""
    (temp_workdir / "notes.txt").write_text("plain text note", encoding="utf-8")
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "notes.txt"})
    assert res.status_code == 200
    body = res.json()
    assert body["language"] == "text"
    assert body["content"] == "plain text note"


@pytest.mark.asyncio
async def test_preview_md_unchanged(temp_workdir, env_cortex_config, reset_deps):
    """回归：.md 仍走 utf-8 直读路径，content 为原文件内容（非 DB 合成）。"""
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "doc1.md"})
    assert res.status_code == 200
    body = res.json()
    assert body["language"] == "markdown"
    # 原文件内容应包含 "Hello world from doc1."
    assert "Hello world from doc1." in body["content"]


@pytest.mark.asyncio
async def test_preview_csv_path_traversal_blocked(env_cortex_config, reset_deps):
    """越权路径应被 _safe_resolve 拦截，返回 FILE_NOT_FOUND。"""
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "../../../etc/passwd.csv"})
    # 路径越权，_safe_resolve 抛 FILE_NOT_FOUND；由于后缀 .csv 在白名单，
    # 走二进制分支前已先 _safe_resolve 校验，返回 FILE_NOT_FOUND 而非 NOT_INDEXED
    assert res.status_code == 404
    assert res.json()["code"] == "FILE_NOT_FOUND"


# ---------------------------------------------------------------------------
# 分页标记（pages 字段）集成测试
# ---------------------------------------------------------------------------


def _init_and_reindex_sync():
    """在非事件循环线程中初始化 + 索引。"""
    idx = deps.get_index_manager()
    idx.reindex(force=True)
    return idx


@pytest.mark.asyncio
async def test_preview_pptx_returns_pages_with_slide_titles(
    temp_workdir, env_cortex_config, reset_deps
):
    """pptx 预览：pages 列出每个 slide，label 含 slide 标题。"""
    from pptx import Presentation  # type: ignore

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Intro Slide"
    s1.placeholders[1].text = "intro body"
    s2 = prs.slides.add_slide(prs.slide_layouts[0])
    s2.shapes.title.text = "Method Slide"
    s2.placeholders[1].text = "method body"
    prs.save(str(temp_workdir / "twoslides.pptx"))

    await asyncio.to_thread(_init_and_reindex_sync)

    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "twoslides.pptx"})

    assert res.status_code == 200
    body = res.json()
    assert body["pages"] is not None
    assert len(body["pages"]) == 2
    assert body["pages"][0]["label"] == "幻灯片 1 · Intro Slide"
    assert body["pages"][1]["label"] == "幻灯片 2 · Method Slide"
    # line_start 是 1-indexed 整数
    assert isinstance(body["pages"][0]["line_start"], int)


@pytest.mark.asyncio
async def test_preview_xlsx_returns_pages_with_sheet_names(
    temp_workdir, env_cortex_config, reset_deps
):
    """xlsx 预览：pages 列出每个 sheet，label 含 sheet 名。"""
    from openpyxl import Workbook  # type: ignore

    wb = Workbook()
    wb.active.title = "Sales"
    wb.active.append(["name", "amt"])
    wb.active.append(["Alice", 100])
    wb.create_sheet("Inventory").append(["item", "qty"])
    wb.save(str(temp_workdir / "multi.xlsx"))

    await asyncio.to_thread(_init_and_reindex_sync)

    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "multi.xlsx"})

    assert res.status_code == 200
    body = res.json()
    assert body["pages"] is not None
    assert len(body["pages"]) == 2
    # _extract_excel_pages 剥掉了 excel_parser 加的 " (N rows)" 后缀
    assert body["pages"][0]["label"] == "工作表 1 · Sales"
    assert body["pages"][1]["label"] == "工作表 2 · Inventory"


@pytest.mark.asyncio
async def test_preview_md_pages_is_null(temp_workdir, env_cortex_config, reset_deps):
    """回归：md 文件预览，pages 为 null（不支持分页）。"""
    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "doc1.md"})

    assert res.status_code == 200
    body = res.json()
    assert body["pages"] is None


@pytest.mark.asyncio
async def test_preview_csv_pages_is_null(temp_workdir, env_cortex_config, reset_deps):
    """回归：csv 文件预览，pages 为 null。"""
    await asyncio.to_thread(_init_and_reindex_sync)

    app = create_app()
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        res = await client.get("/api/preview", params={"path": "data.csv"})

    assert res.status_code == 200
    body = res.json()
    assert body["pages"] is None
