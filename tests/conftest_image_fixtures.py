# -*- coding: utf-8 -*-
"""程序化生成含图 docx/pptx 测试 fixture，避免提交二进制文件。"""
import base64
import io
import os
import tempfile

# 4x4 PNG（PyMuPDF 1.27 拒绝 1x1 PNG："ran out of output before input"）
PNG_FITZ = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAQAAAAECAIAAAAmkwkpAAAAEElEQVR4nGP4z8AARwzEcQCukw/x0F8jngAAAABJRU5ErkJggg=="
)

# 两个不同的 1x1 PNG（不同 sha256），用于去重/多图测试
PNG_RED = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNgAAIAAAUAAeImBZsAAAAASUVORK5CYII="
)
PNG_GREEN = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNgAAIAAAUAAenqB0sAAAAASUVORK5CYII="
)


def make_docx_with_image(path: str, *, with_heading: bool = True) -> str:
    """生成一个含标题 + 段落 + 1 张图片的 docx，返回路径。"""
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    if with_heading:
        doc.add_heading("标题一", level=1)
    doc.add_paragraph("第一段正文。")
    doc.add_picture(io.BytesIO(PNG_RED), width=Inches(1))
    doc.add_paragraph("第二段正文。")
    doc.save(path)
    return path


def make_pptx_with_image(path: str) -> str:
    """生成一个含 1 张标题 slide + 1 张图片的 pptx，返回路径。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.title.text = "幻灯片标题"
    slide.shapes.add_picture(io.BytesIO(PNG_RED), Inches(1), Inches(1), width=Inches(2))
    prs.save(path)
    return path


def make_pdf_with_image(path: str) -> str:
    """生成一个含 1 页文字 + 1 张内嵌 PNG 的 pdf，返回路径。

    使用 ASCII 文字 + fontsize=16：fitz 内置 Helvetica 不支持 CJK，
    pdfplumber 提取会乱码；且两行同字号避免被 _extract_page_text_with_paragraphs
    的字号过滤（main_size*0.85）丢掉第二行。
    """
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 72), "PDF Title Text", fontsize=16)
    page.insert_text((50, 110), "An image should appear below.", fontsize=16)
    page.insert_image(fitz.Rect(50, 130, 200, 280), stream=io.BytesIO(PNG_FITZ))
    doc.save(path)
    doc.close()
    return path
