# -*- coding: utf-8 -*-
"""
@author:XuMing(xuming624@qq.com)
@description: Markitdown-based parser for TreeSearch.

Uses Microsoft's ``markitdown`` library to convert documents to Markdown,
then feeds the result into the existing ``md_to_tree`` pipeline.

Requires optional dependency: ``pip install markitdown``
"""

import logging
import os
import re
from typing import TYPE_CHECKING, Optional

if TYPE_CHECKING:
    from .image_store import ImagePart, ImageStore

logger = logging.getLogger(__name__)

# Extensions handled by markitdown
MARKITDOWN_EXTENSIONS = {".pptx"}

# ---------------------------------------------------------------------------
# Monkey-patch: python-pptx AutoShape.shape_type raises NotImplementedError
# for unrecognized shape types (e.g. SmartArt, connectors, ink).  Returning
# None instead lets markitdown skip those shapes gracefully while still
# extracting text from shapes that have a text_frame.
# ---------------------------------------------------------------------------
def _patch_pptx_shape_type():
    try:
        from pptx.shapes.autoshape import Shape
    except ImportError:
        return

    _orig = Shape.shape_type.fget  # type: ignore[attr-defined]

    @property  # type: ignore[misc]
    def _safe_shape_type(self):
        try:
            return _orig(self)
        except NotImplementedError:
            return None

    Shape.shape_type = _safe_shape_type  # type: ignore[assignment]


_patch_pptx_shape_type()


# ---------------------------------------------------------------------------
# pptx slide-level image extraction + injection helpers
# ---------------------------------------------------------------------------
_RE_SLIDE_COMMENT = re.compile(r"^<!--\s*Slide number:\s*(\d+)\s*-->")
_RE_H1 = re.compile(r"^#\s+\S")


def _extract_pptx_slide_images(pptx_path: str) -> list[list["ImagePart"]]:
    """用 python-pptx 按 slide 顺序提取每 slide 的图片，返回 slide_parts[i]。

    shape_type 取值可能抛 NotImplementedError（见 _patch_pptx_shape_type），
    额外用 getattr 容错。
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from .image_store import ImagePart

    prs = Presentation(pptx_path)
    slide_parts: list[list[ImagePart]] = []
    for idx, slide in enumerate(prs.slides):
        parts: list[ImagePart] = []
        for shape in slide.shapes:
            try:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except (NotImplementedError, AttributeError):
                continue
            try:
                img = shape.image
                disp_w = (
                    round(shape.width / 9525)
                    if getattr(shape, "width", None) else None
                )
                parts.append(
                    ImagePart(
                        blob=img.blob,
                        ext=img.ext or "png",
                        source_ref=f"slide{idx}:{shape.shape_id}",
                        disp_w=disp_w,
                    )
                )
            except Exception as e:  # noqa: BLE001
                logger.warning(
                    "skip pptx picture slide%d shape%s: %s",
                    idx,
                    getattr(shape, "shape_id", "?"),
                    e,
                )
        slide_parts.append(parts)
    return slide_parts


def _inject_slide_images(md: str, slide_mds: list[str]) -> str:
    """按 slide 边界把每 slide 的图片 md 注入到对应 slide 块尾。

    优先按 ``<!-- Slide number: N -->`` 切分；否则按 H1 (``# ``) 切分；
    若切分块数 != slide 数，降级为全部图片追加到末尾。

    注意：comment 和 H1 边界不混合检测 —— markitdown 对同一 slide 同时输出
    ``<!-- Slide number: N -->`` 和 ``# 标题``，混合会重复计数导致注入到
    H1 之前（变成被丢弃的前导文本）。改为分别收集，优先匹配 comment。
    """
    n = len(slide_mds)
    if n == 0 or not any(slide_mds):
        return md

    lines = md.split("\n")

    comment_bounds = [
        i for i, ln in enumerate(lines) if _RE_SLIDE_COMMENT.match(ln.strip())
    ]
    h1_bounds = [
        i for i, ln in enumerate(lines) if _RE_H1.match(ln.strip())
    ]

    boundaries: list[int] = []
    if len(comment_bounds) == n:
        boundaries = comment_bounds
    elif len(h1_bounds) == n:
        boundaries = h1_bounds

    if boundaries:
        # 每个 slide 块尾 = 下一边界行（或文档末尾）；在该位置前插入图片 md
        block_ends = boundaries[1:] + [len(lines)]
        out = list(lines)
        offset = 0
        for i in range(n):
            if not slide_mds[i]:
                continue
            insert_at = block_ends[i] + offset
            block = ["", slide_mds[i], ""]
            out[insert_at:insert_at] = block
            offset += len(block)
        return "\n".join(out)

    # 降级：全部追加末尾
    tail = "\n\n".join(m for m in slide_mds if m)
    return md.rstrip() + "\n\n" + tail


async def markitdown_to_tree(
    file_path: str,
    *,
    model: Optional[str] = None,
    if_add_node_summary: bool = True,
    summary_chars_threshold: int = 600,
    if_add_doc_description: bool = False,
    if_add_node_text: bool = False,
    if_add_node_id: bool = True,
    image_store: "ImageStore | None" = None,
    rel_path: str | None = None,
    **kwargs,
) -> dict:
    """Build a tree index from a document via markitdown.

    Converts the document to Markdown using ``markitdown``, then delegates
    to ``md_to_tree`` for structure extraction.

    Returns:
        {'doc_name': str, 'structure': list, 'source_path': str}
    """
    try:
        from markitdown import MarkItDown
    except ImportError:
        raise ImportError(
            "Markitdown support requires 'markitdown'. "
            "Install with: pip install markitdown"
        )

    doc_name = os.path.splitext(os.path.basename(file_path))[0]
    logger.debug("Parsing with markitdown: %s", file_path)

    md = MarkItDown()
    result = md.convert(file_path)
    md_content = result.text_content

    if not md_content or not md_content.strip():
        logger.warning("markitdown returned empty content for: %s", file_path)
        md_content = ""

    # pptx 图片提取 + slide 级注入（仅 .pptx；其他 markitdown 支持类型暂不提图）
    if image_store is not None and rel_path and file_path.lower().endswith(".pptx"):
        try:
            slide_parts = _extract_pptx_slide_images(file_path)
            all_parts = [p for parts in slide_parts for p in parts]
            refs = (
                image_store.extract_for_doc(rel_path, all_parts)
                if all_parts
                else {}
            )
            slide_mds: list[str] = []
            for idx, parts in enumerate(slide_parts):
                mds = [
                    refs[p.source_ref].inline_md
                    for p in parts
                    if p.source_ref in refs
                ]
                slide_mds.append(" ".join(mds) if mds else "")
            md_content = re.sub(r"!\[[^\]]*\]\(Picture[^)]+\)", "", md_content)
            md_content = _inject_slide_images(md_content, slide_mds)
        except ImportError:
            logger.debug("python-pptx not available, skip pptx image extraction")
        except Exception as e:  # noqa: BLE001
            logger.warning("pptx image extraction failed for %s: %s", file_path, e)

    from ..indexer import md_to_tree

    tree_result = await md_to_tree(
        md_content=md_content,
        model=model,
        if_add_node_summary=if_add_node_summary,
        summary_chars_threshold=summary_chars_threshold,
        if_add_doc_description=if_add_doc_description,
        if_add_node_text=if_add_node_text,
        if_add_node_id=if_add_node_id,
        **kwargs,
    )
    tree_result["doc_name"] = doc_name
    tree_result["source_path"] = os.path.abspath(file_path)
    tree_result["source_type"] = "pptx"

    # PPTX 特殊处理：将扁平的 slide 节点包裹在文档根节点下，形成两层结构
    # markitdown 输出每个 slide 为 # 标题，md_to_tree 将其作为顶层节点
    # 需要聚合为一个根节点 → 多个子节点（每页一个）的层次结构
    structure = tree_result.get("structure", [])
    if len(structure) > 1:
        # 计算整体行号范围
        min_line = min(
            (n.get("line_start", 0) for n in structure if n.get("line_start") is not None),
            default=0,
        )
        max_line = max(
            (n.get("line_end", 0) for n in structure if n.get("line_end") is not None),
            default=0,
        )
        root_node = {
            "title": doc_name,
            "node_id": "0",
            "text": "",
            "line_start": min_line,
            "line_end": max_line,
            "nodes": structure,
        }
        tree_result["structure"] = [root_node]

    return tree_result
